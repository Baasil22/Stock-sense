from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
ppt = Presentation()

# Slide 1: Title Slide
slide_1 = ppt.slides.add_slide(ppt.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "Stock Sense: A Machine Learning Stock Prediction Platform"
subtitle.text = "For IBM Internship\nMohammed Baasil Muneer\nOctober 2024"

# Set title color
title_format = title.text_frame.paragraphs[0].font
title_format.size = Pt(44)
title_format.bold = True
title_format.color.rgb = RGBColor(0, 102, 204)  # Blue color for the title

# Add image in the title slide (e.g., stock market chart or related image)
img_path = "stock_market_chart.jpg"  # Path to your image
slide_1.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), width=Inches(8))

# Slide 2: Introduction with Colors and Bullets
slide_2 = ppt.slides.add_slide(ppt.slide_layouts[1])
title_2 = slide_2.shapes.title
title_2.text = "Introduction"

content_2 = slide_2.shapes.placeholders[1].text_frame
content_2.text = "Project Overview:"
content_2.paragraphs[0].font.size = Pt(28)
content_2.paragraphs[0].font.bold = True
content_2.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)  # Dark blue

p = content_2.add_paragraph()
p.text = "• A platform using machine learning to predict stock market trends."
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(0, 153, 51)  # Green text for emphasis

p = content_2.add_paragraph()
p.text = "• Provides stock analysis, data visualization, and insights."
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(255, 102, 0)  # Orange for a different emphasis

# Slide 3: Problem Statement with Image and Colors
slide_3 = ppt.slides.add_slide(ppt.slide_layouts[1])
title_3 = slide_3.shapes.title
title_3.text = "Problem Statement"

content_3 = slide_3.shapes.placeholders[1].text_frame
content_3.text = "Key Issues:"
content_3.paragraphs[0].font.size = Pt(28)
content_3.paragraphs[0].font.bold = True
content_3.paragraphs[0].font.color.rgb = RGBColor(153, 0, 0)  # Red for problems

p = content_3.add_paragraph()
p.text = "• Uncertainty in stock market trends."
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(204, 0, 0)  # Dark red

p = content_3.add_paragraph()
p.text = "• Difficulty in analyzing large amounts of stock data."
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(204, 0, 0)

# Add an image related to stock issues
img_path_problem = "confused_trader.jpg"  # Path to your image
slide_3.shapes.add_picture(img_path_problem, Inches(5.5), Inches(1.5), width=Inches(3))

# Slide 4: Objectives with Colorful Layout
slide_4 = ppt.slides.add_slide(ppt.slide_layouts[1])
title_4 = slide_4.shapes.title
title_4.text = "Objectives"

content_4 = slide_4.shapes.placeholders[1].text_frame
content_4.text = "Main Goals:"
content_4.paragraphs[0].font.size = Pt(28)
content_4.paragraphs[0].font.bold = True
content_4.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)  # Blue

p = content_4.add_paragraph()
p.text = "• Predict stock prices using machine learning."
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(0, 204, 102)  # Green

p = content_4.add_paragraph()
p.text = "• Provide user-friendly access to stock data."
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(0, 204, 102)

p = content_4.add_paragraph()
p.text = "• Visualize stock market trends."
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(0, 204, 102)

# Slide 5: Technology Stack with Icons
slide_5 = ppt.slides.add_slide(ppt.slide_layouts[1])
title_5 = slide_5.shapes.title
title_5.text = "Technology Stack"

content_5 = slide_5.shapes.placeholders[1].text_frame
content_5.text = "Frontend: HTML5, CSS3, JavaScript"
content_5.paragraphs[0].font.size = Pt(24)
p = content_5.add_paragraph()
p.text = "Backend: Flask"
p.font.size = Pt(24)
p = content_5.add_paragraph()
p.text = "Database: MySQL/SQLite"
p.font.size = Pt(24)
p = content_5.add_paragraph()
p.text = "Data Processing: yfinance API, Scikit-learn"
p.font.size = Pt(24)
p = content_5.add_paragraph()
p.text = "Visualization: Matplotlib, Seaborn"
p.font.size = Pt(24)
p = content_5.add_paragraph()
p.text = "Tools: Streamlit"
p.font.size = Pt(24)

# Add icons to represent each tech stack element (optional)
icon_path = "stack_icon.png"  # Path to a representative tech icon
slide_5.shapes.add_picture(icon_path, Inches(6), Inches(1), width=Inches(1.5))

# Save the presentation
ppt.save("Stock_Sense_IBM_Presentation.pptx")
